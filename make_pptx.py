"""
Génère la présentation PowerPoint complète — Web Sémantique
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG_DARK = RGBColor(0x0F,0x11,0x17)
BG_CARD = RGBColor(0x1E,0x22,0x35)
BLUE    = RGBColor(0x4F,0x8E,0xF7)
PURPLE  = RGBColor(0xA7,0x8B,0xFA)
GREEN   = RGBColor(0x4A,0xDE,0x80)
AMBER   = RGBColor(0xF5,0x9E,0x0B)
RED     = RGBColor(0xF8,0x71,0x71)
TEAL    = RGBColor(0x2D,0xD4,0xBF)
WHITE   = RGBColor(0xFF,0xFF,0xFF)
LIGHT   = RGBColor(0xE2,0xE8,0xF0)
GRAY    = RGBColor(0x64,0x74,0x8B)
W, H    = Inches(13.33), Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG_DARK
    return s

def rect(s, x, y, w, h, color, alpha=None):
    sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh

def circle(s, x, y, size, color):
    sh = s.shapes.add_shape(9, Inches(x), Inches(y), Inches(size), Inches(size))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh

def txbox(s, text, x, y, w, h, size=16, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    return tb

def pill(s, text, x, y, color=BLUE, size=12):
    w = len(text)*0.095 + 0.35
    rect(s, x, y, w, 0.36, color)
    txbox(s, text, x+0.1, y+0.04, w-0.1, 0.3,
          size=size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def header(s, title, accent=BLUE):
    rect(s, 0, 0, 13.33, 0.07, accent)
    txbox(s, title, 0.5, 0.18, 12, 0.65, size=30, bold=True, color=WHITE)
    rect(s, 0.5, 0.83, 1.5, 0.04, accent)

# ═══════════════════════════════════════════
# SLIDE 1 — TITRE
# ═══════════════════════════════════════════
s1 = slide()
rect(s1, 0, 0, 0.1, 7.5, BLUE)
rect(s1, 0.1, 0, 0.05, 7.5, PURPLE)

txbox(s1, "Résumé RDF", 0.5, 1.1, 12, 1.2, size=56, bold=True, color=WHITE)
txbox(s1, "orienté requêtes", 0.5, 2.1, 12, 0.85, size=38, color=BLUE)
txbox(s1, "Une approche basée sur les sous-graphes BGP",
      0.5, 3.1, 11, 0.55, size=19, color=GRAY, italic=True)
rect(s1, 0.5, 3.85, 4.5, 0.04, BLUE)

pill(s1, "Čebirić et al. 2015",  0.5,  4.1, BLUE)
pill(s1, "SumMER 2023",          3.9,  4.1, PURPLE)
pill(s1, "Pappas et al. 2024",   6.6,  4.1, TEAL)
pill(s1, "Notre approche →",     9.5,  4.1, GREEN)

txbox(s1, "État de l'art · Résumé de graphes RDF · Web Sémantique 2025–2026",
      0.5, 6.6, 12, 0.45, size=13, color=GRAY, italic=True)

# ═══════════════════════════════════════════
# SLIDE 2 — PROBLEMATIQUE
# ═══════════════════════════════════════════
s2 = slide()
header(s2, "Problématique", BLUE)

rect(s2, 0.4, 1.1, 5.9, 5.7, BG_CARD)
rect(s2, 0.4, 1.1, 5.9, 0.07, RED)
txbox(s2, "Le problème", 0.65, 1.25, 5.5, 0.45, size=17, bold=True, color=RED)
for txt, yy in [
    ("● Les schémas RDF ont des centaines de classes",         1.85),
    ("● DBpedia → 700+ classes, illisible",                    2.42),
    ("● Les requêtes SPARQL portent sur des motifs structuraux\n  (BGP), pas sur des nœuds isolés", 2.99),
    ("● Les approches existantes traitent les nœuds\n  individuellement → perte du contexte",       3.75),
]:
    txbox(s2, txt, 0.65, yy, 5.4, 0.65, size=13, color=LIGHT)

rect(s2, 6.8, 1.1, 6.1, 5.7, BG_CARD)
rect(s2, 6.8, 1.1, 6.1, 0.07, GREEN)
txbox(s2, "Notre réponse", 7.05, 1.25, 5.8, 0.45, size=17, bold=True, color=GREEN)
for txt, yy in [
    ("✓ Traiter le sous-graphe BGP comme unité atomique",      1.85),
    ("✓ Scorer chaque motif selon l'usage réel",               2.42),
    ("✓ Sélectionner les top-k motifs les plus importants",    2.99),
    ("✓ Connecter via Steiner Tree → résumé cohérent",         3.56),
    ("✓ Évaluer avec Coverage et Compacité",                   4.13),
]:
    txbox(s2, txt, 7.05, yy, 5.8, 0.5, size=13, color=LIGHT)

txbox(s2, "→", 5.9, 3.5, 1.0, 0.7, size=38, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
rect(s2, 0.5, 6.45, 12.4, 0.6, BG_CARD)
txbox(s2, "Idée clé : l'unité naturelle d'usage n'est pas le nœud isolé, mais le sous-graphe BGP.",
      0.5, 6.5, 12.4, 0.5, size=15, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 3 — APPROCHE 1 : Čebirić 2015
# ═══════════════════════════════════════════
s3 = slide()
header(s3, "Approche 1 — Query-Oriented Summarization  (Čebirić et al., 2015)", BLUE)

# Référence
pill(s3, "Section 2 du rapport", 0.5, 0.92, BLUE, 11)

# 3 cartes
for i, (title, body, col) in enumerate([
    ("Objectif",
     "Garantir que le résumé préserve les\ncapacités de réponse aux requêtes SPARQL.\nPas d'importance explicite — approche\nlogique et formelle.",
     BLUE),
    ("Méthode",
     "2 propriétés :\n● Représentativité : toute requête avec\n  réponse sur G en a une sur le résumé\n● Précision : cohérence formelle\n\nBaseline summary O(G²)\nRefined summary O(G⁵)",
     PURPLE),
    ("Évaluation",
     "Exclusivement théorique.\nAucune comparaison quantitative.\nPas de machine learning.\nPas de notion d'importance.",
     AMBER),
]):
    xc = 0.4 + i * 4.3
    rect(s3, xc, 1.25, 4.0, 4.5, BG_CARD)
    rect(s3, xc, 1.25, 4.0, 0.07, col)
    txbox(s3, title, xc+0.2, 1.38, 3.6, 0.45, size=16, bold=True, color=col)
    txbox(s3, body,  xc+0.2, 1.92, 3.6, 3.5, size=13, color=LIGHT)

# Caractéristiques
rect(s3, 0.4, 5.9, 12.5, 1.35, BG_CARD)
txbox(s3, "Caractéristiques clés", 0.65, 5.97, 5, 0.38, size=14, bold=True, color=WHITE)
for txt, xp, col in [
    ("Orientation : Usage/Requêtes",  0.65, BLUE),
    ("Importance explicite : Non",    4.5,  GRAY),
    ("Scalabilité : Faible",          7.8,  RED),
    ("Coverage : Moyenne",           10.5,  AMBER),
]:
    txbox(s3, txt, xp, 6.42, 3.5, 0.38, size=13, color=col)

# ═══════════════════════════════════════════
# SLIDE 4 — APPROCHE 2 : SumMER 2023
# ═══════════════════════════════════════════
s4 = slide()
header(s4, "Approche 2 — SumMER : Structural Summarization  (Troullinou et al., 2023)", PURPLE)

pill(s4, "Section 3 du rapport", 0.5, 0.92, PURPLE, 11)

for i, (title, body, col) in enumerate([
    ("Objectif",
     "Produire un résumé structurel lisible.\nSélectionner les concepts les plus\nimportants via centralité structurelle\net machine learning.",
     PURPLE),
    ("Méthode",
     "① Extraire le schéma RDF\n② 8 mesures de centralité\n   (PageRank, Degree, Betweenness, HITS…)\n③ + nombre d'instances\n④ Régression supervisée\n⑤ Top-k classes → Steiner Tree",
     BLUE),
    ("Évaluation",
     "Empirique et quantitative.\nDatasets réels : DBpedia, SWDF.\nMétrique : Coverage\n→ Améliore significativement\nles approches existantes.",
     GREEN),
]):
    xc = 0.4 + i * 4.3
    rect(s4, xc, 1.25, 4.0, 4.5, BG_CARD)
    rect(s4, xc, 1.25, 4.0, 0.07, col)
    txbox(s4, title, xc+0.2, 1.38, 3.6, 0.45, size=16, bold=True, color=col)
    txbox(s4, body,  xc+0.2, 1.92, 3.6, 3.5, size=13, color=LIGHT)

rect(s4, 0.4, 5.9, 12.5, 1.35, BG_CARD)
txbox(s4, "Caractéristiques clés", 0.65, 5.97, 5, 0.38, size=14, bold=True, color=WHITE)
for txt, xp, col in [
    ("Orientation : Structure",       0.65, PURPLE),
    ("Importance explicite : Oui",    4.5,  GREEN),
    ("Scalabilité : Moyenne",         7.8,  AMBER),
    ("Coverage : Bonne",             10.5,  GREEN),
]:
    txbox(s4, txt, xp, 6.42, 3.5, 0.38, size=13, color=col)

# ═══════════════════════════════════════════
# SLIDE 5 — APPROCHE 3 : Embeddings 2024
# ═══════════════════════════════════════════
s5 = slide()
header(s5, "Approche 3 — Semantic Summaries via Embeddings  (Pappas et al., 2024)", TEAL)

pill(s5, "Section 4 du rapport", 0.5, 0.92, TEAL, 11)

for i, (title, body, col) in enumerate([
    ("Objectif",
     "Améliorer scalabilité et qualité\nsémantique. Remplacer les mesures de\ncentralité explicites par des\nreprésentations vectorielles apprises\nautomatiquement (RDF2Vec).",
     TEAL),
    ("Méthode",
     "① Random walks sur le graphe RDF\n② Word2Vec → embeddings ∈ ℝᵈ\n③ Labels d'importance depuis logs\n   I_node(v) = requêtes(v)/total\n④ Decision Tree (nœuds)\n   SVR (arêtes)\n⑤ Top-k + SDIST (Steiner pondéré)",
     BLUE),
    ("Évaluation",
     "Quantitative et comparative.\nComparée à SumMER (référence).\nMétriques : Coverage,\nvrais positifs, temps d'exécution.\n→ Meilleure coverage et scalabilité.",
     GREEN),
]):
    xc = 0.4 + i * 4.3
    rect(s5, xc, 1.25, 4.0, 4.5, BG_CARD)
    rect(s5, xc, 1.25, 4.0, 0.07, col)
    txbox(s5, title, xc+0.2, 1.38, 3.6, 0.45, size=16, bold=True, color=col)
    txbox(s5, body,  xc+0.2, 1.92, 3.6, 3.5, size=13, color=LIGHT)

rect(s5, 0.4, 5.9, 12.5, 1.35, BG_CARD)
txbox(s5, "Caractéristiques clés", 0.65, 5.97, 5, 0.38, size=14, bold=True, color=WHITE)
for txt, xp, col in [
    ("Orientation : Sémantique + Usage",  0.65, TEAL),
    ("Importance : Apprise (ML)",         4.5,  GREEN),
    ("Scalabilité : Élevée",              7.8,  GREEN),
    ("Coverage : Meilleure",             10.5,  GREEN),
]:
    txbox(s5, txt, xp, 6.42, 3.5, 0.38, size=13, color=col)

# ═══════════════════════════════════════════
# SLIDE 6 — TABLEAU COMPARATIF
# ═══════════════════════════════════════════
s6 = slide()
header(s6, "Synthèse comparative des 3 approches  (Section 5)", AMBER)

# Colonnes header
cols = ["Critère", "Čebirić 2015", "SumMER 2023", "Pappas 2024"]
col_colors = [GRAY, BLUE, PURPLE, TEAL]
col_x = [0.4, 3.2, 6.1, 9.0]
col_w = [2.7, 2.8, 2.8, 4.0]

for j, (label, col, xp, ww) in enumerate(zip(cols, col_colors, col_x, col_w)):
    rect(s6, xp, 1.1, ww, 0.55, col)
    txbox(s6, label, xp+0.1, 1.17, ww-0.1, 0.42,
          size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("Orientation",          "Usage",           "Structure",        "Sémantique+Usage"),
    ("Unité du résumé",      "BGP (motifs)",    "Nœuds isolés",     "Nœuds isolés"),
    ("Importance explicite", "Non",             "Oui (centralités)","Oui (apprise)"),
    ("Machine Learning",     "Non",             "Oui (régression)", "Oui (embeddings)"),
    ("Pondération arêtes",   "Non",             "Non",              "Oui (SDIST)"),
    ("Scalabilité",          "Faible",          "Moyenne",          "Élevée"),
    ("Coverage",             "Moyenne",         "Bonne",            "Meilleure"),
    ("Interprétabilité",     "Élevée",          "Élevée",           "Plus limitée"),
]
val_colors = {
    "Élevée": GREEN, "Meilleure": GREEN, "Bonne": GREEN,
    "Oui": GREEN, "Oui (centralités)": GREEN, "Oui (apprise)": GREEN,
    "Oui (régression)": GREEN, "Oui (embeddings)": GREEN, "Oui (SDIST)": GREEN,
    "Faible": RED, "Non": GRAY, "Moyenne": AMBER,
    "Plus limitée": AMBER,
}

for i, row in enumerate(rows):
    yy = 1.75 + i * 0.58
    bg = BG_CARD if i % 2 == 0 else RGBColor(0x16,0x19,0x29)
    for j, (val, xp, ww) in enumerate(zip(row, col_x, col_w)):
        rect(s6, xp, yy, ww, 0.52, bg)
        col = WHITE if j == 0 else val_colors.get(val, LIGHT)
        txbox(s6, val, xp+0.1, yy+0.08, ww-0.15, 0.38,
              size=12, bold=(j==0), color=col, align=PP_ALIGN.CENTER)

rect(s6, 0.4, 6.45, 12.5, 0.75, RGBColor(0x1a,0x3a,0x2a))
txbox(s6, "→  Lacune identifiée : toutes traitent des éléments isolés — les requêtes SPARQL\n"
          "    expriment des motifs BGP impliquant plusieurs éléments simultanément.",
      0.6, 6.5, 12.2, 0.65, size=13, bold=True, color=GREEN)

# ═══════════════════════════════════════════
# SLIDE 7 — DATASET
# ═══════════════════════════════════════════
s7 = slide()
header(s7, "Dataset — SWDF (Semantic Web Dog Food)", GREEN)

for i, (icon, title, body, col) in enumerate([
    ("📄", "Schéma RDF",
     "Ontologie SWC (2009-05-11)\n18 classes · 13 propriétés\n\nSource :\nlov.linkeddata.es",
     BLUE),
    ("🔍", "Requêtes SPARQL",
     "15 requêtes dont\n3 vraies requêtes LSQ\n(logs réels mai 2014)\n\nSource : GitHub/AKSW",
     PURPLE),
    ("📋", "Logs Apache",
     "2 006 lignes de vrais logs\ndu serveur\ndata.semanticweb.org\n\nSource : LSQ dataset",
     GREEN),
]):
    xc = 0.4 + i * 4.3
    rect(s7, xc, 1.1, 4.0, 5.6, BG_CARD)
    rect(s7, xc, 1.1, 4.0, 0.07, col)
    txbox(s7, icon,  xc+0.2, 1.25, 0.7, 0.6, size=28)
    txbox(s7, title, xc+0.2, 1.85, 3.6, 0.45, size=16, bold=True, color=col)
    txbox(s7, body,  xc+0.2, 2.42, 3.6, 3.8, size=13, color=LIGHT)

rect(s7, 0.4, 6.85, 12.5, 0.5, BG_CARD)
txbox(s7, "Données réelles — ontologie + logs téléchargés depuis les sources officielles",
      0.4, 6.9, 12.5, 0.4, size=13, bold=True, color=AMBER, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════
# SLIDE 8 — PIPELINE
# ═══════════════════════════════════════════
s8 = slide()
header(s8, "Notre approche — Pipeline en 5 étapes  (Section 7)", AMBER)

steps = [
    (1, "Parser BGP",    "Extraire classes\net propriétés\nde chaque requête\nSPARQL", BLUE),
    (2, "Sous-graphes",  "Construire S_q :\nportion du schéma\ninduite par\nle BGP",  PURPLE),
    (3, "Scorer I(S)",   "Calculer le score\nd'importance\npar normalisation\nnaturelle", GREEN),
    (4, "Top-k",         "Sélectionner\nles k meilleurs\nsous-graphes",                AMBER),
    (5, "Steiner Tree",  "Connecter les\nsous-graphes →\nrésumé connexe\nfinal",       RED),
]
for i, (n, t, d, c) in enumerate(steps):
    xc = 0.35 + i * 2.58
    rect(s8, xc, 1.1, 2.35, 3.3, BG_CARD)
    rect(s8, xc, 1.1, 2.35, 0.07, c)
    cc = circle(s8, xc+0.9, 1.22, 0.5, c)
    tf = cc.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = str(n)
    run.font.size = Pt(14); run.font.bold = True; run.font.color.rgb = WHITE
    txbox(s8, t, xc+0.1, 1.82, 2.1, 0.45, size=13, bold=True, color=c)
    txbox(s8, d, xc+0.1, 2.35, 2.1, 1.5,  size=12, color=LIGHT)
    if i < 4:
        txbox(s8, "→", 2.55+i*2.58, 2.2, 0.35, 0.5,
              size=20, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

# Formule + normalisation
rect(s8, 0.4, 4.6, 12.5, 0.65, BG_CARD)
txbox(s8, "I(S)  =  α · freq_n(S)  +  β · EC_n(S)  +  γ · Div_n(S)",
      0.4, 4.66, 12.5, 0.55,
      size=21, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 3 définitions de normalisation
norm_items = [
    ("freq_n(S)",  "freq(S) / total_requêtes",      "déjà dans [0, 1]",  BLUE,   1.0),
    ("EC_n(S)",    "EC(S) / |E(schéma)|",           "EC brut / 13 arêtes", PURPLE, 4.5),
    ("Div_n(S)",   "|props distinctes| / |E(S)|",   "déjà dans [0, 1]",  GREEN,  8.1),
]
for sym, formula, note, col, xp in norm_items:
    rect(s8, xp, 5.42, 4.0, 1.82, BG_CARD)
    txbox(s8, sym,     xp+0.15, 5.5,  3.7, 0.38, size=15, bold=True, color=col)
    txbox(s8, formula, xp+0.15, 5.9,  3.7, 0.38, size=13, color=LIGHT)
    txbox(s8, note,    xp+0.15, 6.28, 3.7, 0.38, size=11, color=GRAY, italic=True)

txbox(s8, "α = 0.5  |  β = 0.3  |  γ = 0.2   avec α + β + γ = 1",
      0.4, 7.1, 12.5, 0.35,
      size=13, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════
# SLIDE 9 — RÉSULTATS
# ═══════════════════════════════════════════
s9 = slide()
header(s9, "Résultats — Dataset SWDF  (k = 5, α=0.5, β=0.3, γ=0.2)", BLUE)

# 4 métriques
for val, lbl, xp, col in [
    ("50%",  "Coverage globale",  0.4,  GREEN),
    ("58%",  "Coverage nœuds",    3.55, BLUE),
    ("37%",  "Coverage arêtes",   6.7,  PURPLE),
    ("61%",  "Compacité",         9.85, AMBER),
]:
    rect(s9, xp, 1.05, 3.0, 1.55, BG_CARD)
    txbox(s9, val, xp, 1.12, 3.0, 0.82,
          size=38, bold=True, color=col, align=PP_ALIGN.CENTER)
    txbox(s9, lbl, xp, 1.88, 3.0, 0.38,
          size=13, color=GRAY, align=PP_ALIGN.CENTER)

# Tableau k vs coverage
rect(s9, 0.4, 2.8, 6.3, 4.45, BG_CARD)
txbox(s9, "Impact de k", 0.6, 2.9, 6.0, 0.42, size=15, bold=True, color=WHITE)
rect(s9, 0.45, 3.38, 6.2, 0.04, BLUE)
for hdr, xp in [("k", 0.6), ("Coverage", 1.6), ("Compacité", 3.6), ("Remarque", 5.1)]:
    txbox(s9, hdr, xp, 3.04, 1.8, 0.35, size=12, bold=True, color=BLUE)

rows_k = [
    ("1",  "25%",  "77%", "Résumé minimal"),
    ("5",  "50%",  "61%", "✓ Meilleur compromis"),
    ("7",  "85%",  "35%", "Plus complet"),
    ("14", "100%", "26%", "= schéma entier"),
]
for i, (k_v, cov, comp, rem) in enumerate(rows_k):
    yy = 3.5 + i * 0.82
    is_best = (i == 1)
    if is_best:
        rect(s9, 0.45, yy-0.06, 6.2, 0.75, RGBColor(0x1a,0x3a,0x2a))
    col_row = GREEN if is_best else LIGHT
    for val, xp in [(k_v, 0.6), (cov, 1.6), (comp, 3.6), (rem, 5.1)]:
        txbox(s9, val, xp, yy, 1.8, 0.55, size=13,
              bold=is_best, color=col_row)

# Interprétation
rect(s9, 7.0, 2.8, 6.0, 4.45, BG_CARD)
txbox(s9, "Interprétation", 7.2, 2.9, 5.7, 0.42, size=15, bold=True, color=WHITE)
for txt, yy, col in [
    ("Coverage 50% :", 3.5, BLUE),
    ("La moitié des éléments utilisés\ndans les requêtes est présente\ndans le résumé.", 3.88, LIGHT),
    ("Compacité 61% :", 4.85, AMBER),
    ("Le résumé = seulement 39%\nde la taille du schéma original.", 5.23, LIGHT),
    ("→ Bon compromis couverture / taille", 6.1, GREEN),
]:
    txbox(s9, txt, 7.2, yy, 5.5, 0.7, size=13, bold=(col!=LIGHT), color=col)

# ═══════════════════════════════════════════
# SLIDE 10 — DEMO
# ═══════════════════════════════════════════
s10 = slide()
header(s10, "Démonstration — Interface Streamlit", PURPLE)

for i, (title, body, col) in enumerate([
    ("📊  Schéma RDF",    "Graphe interactif\n18 classes · 13 propriétés\nZoom & hover",          BLUE),
    ("⚙️  Pipeline",      "Lancement en 1 clic\nTableau des scores I(S)\nfréq_n · EC_n · Div_n", PURPLE),
    ("📈  Évaluation",    "Jauges Coverage\net Compacité\nÉléments manquants",                    GREEN),
    ("🔬  Analyse",       "Coverage vs k\nImpact des poids\nTableaux comparatifs",               AMBER),
]):
    xc = 0.4 + i * 3.2
    rect(s10, xc, 1.1, 3.0, 4.1, BG_CARD)
    rect(s10, xc, 1.1, 3.0, 0.07, col)
    txbox(s10, title, xc+0.15, 1.25, 2.7, 0.5, size=14, bold=True, color=col)
    txbox(s10, body,  xc+0.15, 1.85, 2.7, 2.9, size=13, color=LIGHT)

rect(s10, 0.4, 5.35, 12.5, 1.95, BG_CARD)
txbox(s10, "Paramètres ajustables en temps réel", 0.6, 5.43, 8, 0.42, size=14, bold=True, color=WHITE)
for sym, lbl, rng, col, xp in [
    ("k",  "Nombre de sous-graphes", "1 → 14",  BLUE,   0.6),
    ("α",  "Poids fréquence",        "0 → 1",   PURPLE, 3.7),
    ("β",  "Poids connectivité",     "0 → 1-α", GREEN,  6.8),
    ("γ",  "= 1 - α - β",           "auto",    AMBER,  9.9),
]:
    txbox(s10, sym, xp, 5.95, 0.5, 0.55, size=24, bold=True, color=col)
    txbox(s10, lbl, xp+0.5, 6.0,  2.7, 0.32, size=12, color=LIGHT)
    txbox(s10, rng, xp+0.5, 6.38, 2.7, 0.32, size=11, color=GRAY)

txbox(s10, "streamlit run app.py   →   http://localhost:8501",
      0.4, 7.1, 12.5, 0.35, size=13, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════
# SLIDE 11 — CONCLUSION
# ═══════════════════════════════════════════
s11 = slide()
rect(s11, 0, 0, 0.1, 7.5, BLUE)
rect(s11, 0.1, 0, 0.05, 7.5, PURPLE)

txbox(s11, "Conclusion", 0.5, 0.35, 12, 0.75, size=38, bold=True, color=WHITE)
rect(s11, 0.5, 1.1, 2, 0.05, BLUE)

for col, sym, title, body, yy in [
    (GREEN,  "✓", "Approche orientée sous-graphes BGP",
     "Contrairement aux 3 approches étudiées qui traitent des nœuds isolés,\nnous traitons les motifs de requêtes comme unité atomique.", 1.25),
    (BLUE,   "✓", "Normalisation cohérente",
     "freq_n = freq/total_q  ·  EC_n = EC/|E(schéma)|  ·  Div_n déjà dans [0,1]", 2.4),
    (PURPLE, "✓", "Données réelles SWDF + LSQ",
     "Ontologie SWC officielle + logs Apache réels du serveur data.semanticweb.org", 3.3),
    (AMBER,  "✓", "Interface interactive Streamlit",
     "Paramètres ajustables, graphes Plotly interactifs, analyse comparative", 4.2),
]:
    cc = circle(s11, 0.5, yy+0.05, 0.5, col)
    tf = cc.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = sym
    run.font.size = Pt(13); run.font.bold = True; run.font.color.rgb = WHITE
    txbox(s11, title, 1.2, yy+0.05, 11.5, 0.4, size=16, bold=True, color=col)
    txbox(s11, body,  1.2, yy+0.47, 11.5, 0.65, size=13, color=LIGHT)

rect(s11, 0.5, 5.4, 12.3, 0.75, BG_CARD)
txbox(s11, "k = 5  →  Coverage 50%  avec Compacité 61%  (résumé = 39% du schéma original)",
      0.5, 5.52, 12.3, 0.5, size=17, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

rect(s11, 0.5, 6.35, 12.3, 0.9, BG_CARD)
txbox(s11,
      "Perspective : avec un dataset plus grand (milliers de requêtes), freq(S) varierait\n"
      "davantage → meilleure discrimination entre sous-graphes.",
      0.65, 6.42, 12.0, 0.75, size=13, color=GRAY, italic=True)

# ═══════════════════════════════════════════
# SLIDE 12 — OUTILS (bref)
# ═══════════════════════════════════════════
s12 = slide()
header(s12, "Outils utilisés", GRAY)

tools = [
    ("Python 3.14",   "Langage principal",              BLUE),
    ("rdflib 7.6",    "Chargement schéma RDF/TTL",      BLUE),
    ("NetworkX 3.6",  "Graphes + Steiner Tree",          BLUE),
    ("SPARQLWrapper", "Parsing requêtes SPARQL",         PURPLE),
    ("Streamlit",     "Interface web interactive",       GREEN),
    ("Plotly 6.6",    "Graphiques interactifs",           GREEN),
    ("Jupyter",       "Notebook de démonstration",       AMBER),
    ("python-pptx",   "Génération de ce document",       AMBER),
]
for i, (tool, desc, col) in enumerate(tools):
    cn = i // 4; rn = i % 4
    xp = 0.4 + cn * 6.4; yp = 1.1 + rn * 1.5
    rect(s12, xp, yp, 6.0, 1.3, BG_CARD)
    rect(s12, xp, yp, 6.0, 0.06, col)
    txbox(s12, tool, xp+0.2, yp+0.15, 5.6, 0.42, size=16, bold=True, color=col)
    txbox(s12, desc, xp+0.2, yp+0.6,  5.6, 0.42, size=13, color=GRAY)

txbox(s12, "Toutes les bibliothèques sont open-source",
      0.4, 7.05, 12.5, 0.35, size=12, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════
out = "/Users/a-m-s-c/websemantique_presentation.pptx"
prs.save(out)
print(f"✓ {out}")
print(f"  12 slides générées")
